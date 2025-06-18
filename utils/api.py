# Copyright 2024 DataRobot, Inc.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#   http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
from __future__ import annotations
import io # ファイル処理用
import base64 # ファイルエンコーディング用
import os # ファイルパス処理用

import ast
import asyncio
import inspect
import json
import logging
from joblib import Memory
from copy import deepcopy
from datetime import datetime

import re
import sys
import tempfile
from dataclasses import dataclass
from types import  TracebackType
from typing import (
    Any,
    AsyncGenerator,
    Type,
    TypeVar,
    cast,
)

import datarobot as dr
import instructor
import numpy as np
import pandas as pd
import psutil
from datarobot.client import RESTClientObject
import openai
from openai import AsyncOpenAI
from openai.types.chat.chat_completion_message_param import ChatCompletionMessageParam
from openai.types.chat.chat_completion_system_message_param import (
    ChatCompletionSystemMessageParam,
)
import streamlit as st


# Office文書読み込み用のライブラリ
from docx import Document  # Word文書処理用
from pptx import Presentation  # PowerPoint処理用

from pydantic import ValidationError

sys.path.append("..")
from utils import prompts
from utils.logging_helper import get_logger, log_api_call
from utils.dr_helper import (
    initialize_deployment,
    fetch_aicatalog_dataset,
    async_submit_actuals_to_datarobot
    )
from utils.schema import PromptType

logger = get_logger()
logging.getLogger("openai").setLevel(logging.WARNING)
logging.getLogger("openai.http_client").setLevel(logging.WARNING)

def log_memory() -> None:
    process = psutil.Process()
    memory = process.memory_info().rss / 1024 / 1024  # MB
    logger.info(f"Memory usage: {memory:.2f} MB")



class AsyncLLMClient:

    def __init__(self, use_instructor:bool=False) -> None:
        self.use_instructor = use_instructor

    async def __aenter__(self) -> instructor.AsyncInstructor:
        dr_client, deployment_chat_base_url = initialize_deployment()
        
        self.openai_client = AsyncOpenAI(
            api_key=dr_client.token,
            base_url=deployment_chat_base_url,
            timeout=90,
            max_retries=2,
        )

        if self.use_instructor:
            self.client = instructor.from_openai(
                self.openai_client, mode=instructor.Mode.MD_JSON
            )
        else:
            self.client = self.openai_client
        return self.client

    async def __aexit__(
        self,
        exc_type: Type[BaseException] | None,
        exc_val: BaseException | None,
        exc_tb: TracebackType | None,
    ) -> None:
        await self.openai_client.close()  # Properly close the client


ALTERNATIVE_LLM_BIG = "datarobot-deployed-llm"
ALTERNATIVE_LLM_SMALL = "datarobot-deployed-llm"
DICTIONARY_BATCH_SIZE = 10
MAX_REGISTRY_DATASET_SIZE = 400e6  # aligns to 400MB set in streamlit config.toml
DISK_CACHE_LIMIT_BYTES = 512e6

_memory = Memory(tempfile.gettempdir(), verbose=0)
_memory.clear(warn=False)  # clear cache on startup

T = TypeVar("T")


def cache(f: T) -> T:
    """Cache function and coroutine results to disk using joblib."""
    cached_f = _memory.cache(f)

    if asyncio.iscoroutinefunction(f):

        async def awrapper(*args: Any, **kwargs: Any) -> Any:
            in_cache = cached_f.check_call_in_cache(*args, **kwargs)
            result = await cached_f(*args, **kwargs)
            if not in_cache:
                _memory.reduce_size(DISK_CACHE_LIMIT_BYTES)
            else:
                logger.info(
                    f"Using previously cached result for function `{f.__name__}`"
                )
            return result

        return cast(T, awrapper)
    else:

        def wrapper(*args: Any, **kwargs: Any) -> Any:
            in_cache = cached_f.check_call_in_cache(*args, **kwargs)
            result = cached_f(*args, **kwargs)
            if not in_cache:
                _memory.reduce_size(DISK_CACHE_LIMIT_BYTES)
            else:
                logger.info(
                    f"Using previously cached result for function `{f.__name__}`"  # type: ignore[attr-defined]
                )
            return result

        return cast(T, wrapper)

# --- ファイルアップロード処理関数 ---
def process_uploaded_file(uploaded_file):
    """アップロードされたファイルを処理し、内容を抽出します。"""
    file_type = uploaded_file.type
    file_name = uploaded_file.name
    file_content = None
    file_summary = None

    try:
        # ファイルタイプに応じた処理
        if file_type == "text/csv":
            # CSVファイルの処理
            df = pd.read_csv(uploaded_file)
            file_content = df.head(10).to_dict()
            file_summary = f"CSVファイル: {len(df.index)}行 × {len(df.columns)}列。列名: {', '.join(df.columns.tolist())}"

        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or \
            file_type == "application/vnd.ms-excel":
            # Excelファイルの処理
            df = pd.read_excel(uploaded_file)
            file_content = df.head(10).to_dict()  # 最初の10行を辞書形式で取得
            file_summary = f"Excelファイル: {len(df.index)}行 × {len(df.columns)}列。列名: {', '.join(df.columns.tolist())}"

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # Word文書の処理
            doc = Document(uploaded_file)
            text_content = []
            for para in doc.paragraphs:
                if para.text.strip():  # 空のパラグラフはスキップ
                    text_content.append(para.text)
            file_content = "\n".join(text_content)
            file_summary = f"Word文書: {len(text_content)}段落。先頭100文字: {file_content[:100]}..."

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # PowerPointの処理
            prs = Presentation(uploaded_file)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
            file_content = " ".join(text_content)
            file_summary = f"PowerPointファイル: {len(prs.slides)}スライド。テキスト先頭100文字: {file_content[:100]}..."

        else:
            file_content = "サポートされていないファイル形式です。"
            file_summary = f"サポートされていないファイル形式: {file_type}"

        return {
            "content_type": file_type,
            "content": file_content,
            "summary": file_summary
        }

    except Exception as e:
        return {
            "content_type": file_type,
            "content": f"ファイル処理中にエラーが発生しました: {str(e)}",
            "summary": f"処理エラー: {str(e)}"
        }

async def fetch_prompts_with_tools(use_tools_and_descriptions:bool=True, 
                                   prompt_type:str=None) -> str:
    if use_tools_and_descriptions:
        system_prompt = system_prompt_switcher(prompt_type=prompt_type)
        logger.info(f"Use descriptions")
        df = await fetch_aicatalog_dataset()
        _data = df.to_dict(orient='list')
        print(_data)
        tools_and_descriptions = "\n".join([str(i)+':'+str(j) for i, j in zip(_data["tool_name"], _data['description'])])
        if prompt_type is None:
           system_prompt = system_prompt.format(tools_and_descriptions=tools_and_descriptions,
                                                current_question_round=st.session_state.question_counter)
        else:
            system_prompt = system_prompt.format(tools_and_descriptions=tools_and_descriptions)
        logger.info(f"system_prompt: {system_prompt}")
    else:
        logger.info(f"Don't use descriptions")
        system_prompt = prompts.get_system_prompt()
        tools = prompts.DX_TOOLS
        tools = ", ".join(tools)
        system_prompt = system_prompt.format(tools=tools,
                                             current_question_round=st.session_state.question_counter)
        
    return system_prompt

async def prepare_telemetry_send(telemetry_json: dict | None) -> dict | None:
    """telemetry_jsonをdeepcopyし、startTimestampを追加して返す"""
    if not telemetry_json:
        return None
    
    telemetry_send = deepcopy(telemetry_json)
    telemetry_send["startTimestamp"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    return telemetry_send

def system_prompt_switcher(prompt_type:str=None) -> str:
    """
    システムプロンプトを取得します。
    セッション状態に基づいて決定されます。
    prompt_typeがNoneの場合は全て同じシステムプロンプトで処理を行います。
    prompt_typeが"decision"の場合は決定支援のプロンプトを使用します。
    prompt_typeが"question"の場合は質問作成のプロンプトを使用します。
    prompt_typeが"solution"の場合はソリューション提案のプロンプトを使用します。
    """
    if prompt_type is None:
        logger.info(f"Use default system prompt")
        return prompts.get_system_prompt_description()
    elif prompt_type == PromptType.DECISION:
        # 決定支援のためのプロンプトを取得
        logger.info(f"Use decision support prompt")
        return prompts.get_system_prompt_for_decision()
    elif prompt_type == PromptType.QUESTION:
        # 質問応答のためのプロンプトを取得
        logger.info(f"Use question answering prompt")
        return prompts.get_system_prompt_for_questions()
    elif prompt_type == PromptType.SOLUTION:
        # ソリューション提案のためのプロンプトを取得
        logger.info(f"Use solution proposal prompt")
        return prompts.get_system_prompt_for_solution()


async def get_llm_responses(messages:str,
                            use_instructor:bool=False,
                            response_model:str=None) -> tuple[str, str]:
    if use_instructor:
        async with AsyncLLMClient(use_instructor=use_instructor) as client:
            (
            completion,
            completion_org,
            ) = await client.chat.completions.create_with_completion(
            response_model=response_model,
            model=ALTERNATIVE_LLM_BIG,
            temperature=0.1,
            messages=messages,
            )
        # ここのvalidationの処理はpydanticを使用してやりたいが分岐が多いので一旦動くのかを考える

        response_content = completion
        association_id = completion_org.datarobot_moderations["association_id"]

    else:
        async with AsyncLLMClient(use_instructor=use_instructor) as client:
            completion = await client.chat.completions.create(
                response_format={"type": "json_object"},
                model=ALTERNATIVE_LLM_BIG,
                messages=messages,
            )
        response_content = completion.choices[0].message.content
        association_id = completion.datarobot_moderations['association_id']

    logger.debug("raw LLM response:" +completion)
    return response_content, association_id

@log_api_call
# --- API 呼び出し関数 ---
async def fetch_dx_tool_suggestions(chat_history_for_openai,
                                    use_tools_and_descriptions:bool=True,
                                    telemetry_json:dict | None =None,
                                    prompt_type:str = None
                                    ) -> dict:
    """
    OpenAI APIを呼び出して、DXテーマ定義に関する応答を取得します。
    AIにはJSON形式で応答するよう指示します。
    """
    # use_instructor = True if prompt_type is not None else False

    system_prompt = await fetch_prompts_with_tools(use_tools_and_descriptions,
                                                   prompt_type)

    messages: list[ChatCompletionMessageParam] = [
            ChatCompletionSystemMessageParam(
                role="system", content=system_prompt
            )]

    messages.extend(chat_history_for_openai) # これまでの会話履歴を追加
    print(messages)

    #　送信するデータのハンドリング

    telemetry_send = await prepare_telemetry_send(telemetry_json)

    try:

        async with AsyncLLMClient() as client:
                #(
            #    completion,
            #    completion_org,
            #) = await client.chat.completions.create_with_completion(
            completion = await client.chat.completions.create(
                response_format={"type": "json_object"},
                model="datarobot-deployed-llm",
                messages=messages,
            )
        print(completion)
        response_content = completion.choices[0].message.content
        association_id = completion.datarobot_moderations['association_id']
        logger.info(f"Association ID: {association_id}")
        logger.info(f"telemetry_send: {telemetry_send}")
        logger.info(f"telemetry_send is truthy: {bool(telemetry_send)}")

        if telemetry_send:
            logger.info("submit telemetry")
            #task = asyncio.create_task(
            await async_submit_actuals_to_datarobot(
                    association_id=association_id, telemetry_json=telemetry_send
                )
            #)

        # デバッグ用にAIの生の応答をコンソールに出力
        print("--- OpenAI Raw Response ---")
        print(response_content)
        print("--------------------------")

        parsed_response = json.loads(response_content)
        # 応答のバリデーション。pydanticを使用していないため、手動でチェック。後で書き換える。
        if not isinstance(parsed_response, dict) or "type" not in parsed_response or "message" not in parsed_response:
            raise ValueError("AIの応答に必要な'type'または'message'フィールドが含まれていません。")
        if parsed_response["type"] == "solution":
            if "tools" not in parsed_response and "tool" not in parsed_response:
                # 後方互換性のために単一ツールの場合も処理
                raise ValueError("AIの'solution'タイプの応答に必要な'tools'または'tool'フィールドが含まれていません。")
            if "todos" not in parsed_response:
                raise ValueError("AIの'solution'タイプの応答に必要な'todos'フィールドが含まれていません。")
            # 旧形式の応答を新形式に変換（後方互換性のため）
            if "tool" in parsed_response and "tools" not in parsed_response:
                parsed_response["tools"] = [parsed_response["tool"]]
                parsed_response["primary_tool"] = parsed_response["tool"]
                if "tool_combinations" not in parsed_response:
                    parsed_response["tool_combinations"] = [{
                        "tool": parsed_response["tool"],
                        "purpose": "主要な解決手段",
                        "todos": parsed_response["todos"]
                    }]
        if parsed_response["type"] == "questions" and "questions" not in parsed_response:
            raise ValueError("AIの'questions'タイプの応答に必要な'questions'フィールドが含まれていません。")

        return parsed_response

    except openai.APIError as e:
        st.error(f"OpenAI APIエラーが発生しました: {e}")
        print(f"OpenAI API Error: {e}")
        return {"type": "error", "message": f"APIとの通信に失敗しました。詳細: {str(e)}"}
    except json.JSONDecodeError:
        st.error(f"AIの応答の解析に失敗しました。AIが有効なJSONを返さなかった可能性があります。Raw response: {response_content}")
        print(f"JSONDecodeError. Raw response: {response_content}")
        return {"type": "error", "message": "AIの応答形式が正しくありません。"}
    except ValueError as e: # 自作のバリデーションエラー
        st.error(f"AIの応答構造が不正です: {e} Raw response: {response_content}")
        print(f"ValueError (response structure). Raw response: {response_content}")
        return {"type": "error", "message": f"AIの応答構造が予期されたものではありません。詳細: {str(e)}"}
    except Exception as e:
        st.error(f"予期せぬエラーが発生しました: {e}")
        print(f"Unexpected error: {e}")
        return {"type": "error", "message": f"処理中に予期せぬエラーが発生しました。詳細: {str(e)}"}



    # --- ファイルのAI要約処理 ---
def summarize_file_content(file_info, filename):
    """AIを使用してファイルの内容を要約します。"""
    # ここではシンプルな要約を返しますが、将来的にはAI要約機能を実装できます
    return file_info["summary"]



